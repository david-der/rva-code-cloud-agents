"""
PowerPoint Agent for creating presentations with AI-generated backgrounds
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os
import io
import requests
from typing import Literal, List, Dict
from .base_agent import BaseAgent
import time
import traceback

class PowerPointAgent(BaseAgent):
    def generate(self,
                prompt: str,
                slides_config: List[Dict] = None,
                output_dir: str = "output",
                timeout: int = 30) -> dict:
        """
        Generate a PowerPoint presentation with AI-generated background and charts
        
        Args:
            prompt: The image generation prompt for background
            slides_config: List of slide configurations
                Example: [
                    {
                        "type": "title",
                        "text": "AI in Industrial Analytics",
                        "text_color": "white"
                    },
                    {
                        "type": "chart",
                        "chart_path": "path/to/chart.png"
                    }
                ]
            output_dir: Directory to save outputs
            timeout: Request timeout in seconds
        """
        start_time = time.time()
        print("Starting PowerPoint generation...")
        
        # Default slides config if none provided
        if slides_config is None:
            slides_config = [
                {
                    "type": "title",
                    "text": "AI Generated Presentation",
                    "text_color": "white"
                }
            ]
        
        try:
            os.makedirs(output_dir, exist_ok=True)
            
            # Set up paths
            image_path = os.path.join(output_dir, "generated_image.png")
            pptx_path = os.path.join(output_dir, "presentation.pptx")
            
            # Generate or reuse background image
            if not os.path.exists(image_path):
                print("Generating new background image...")
                response = self.client.images.generate(
                    model="dall-e-3",
                    prompt=prompt,
                    size="1792x1024",
                    quality="standard",
                    n=1,
                )
                image_url = response.data[0].url
                image_response = requests.get(image_url, timeout=timeout)
                image_response.raise_for_status()
                image = Image.open(io.BytesIO(image_response.content))
                image.save(image_path)
            else:
                print("Using existing background image...")
            
            # Create presentation
            prs = Presentation()
            
            # Process each slide
            for slide_config in slides_config:
                if slide_config["type"] == "title":
                    self._create_slide_with_image(
                        prs,
                        image_path,
                        slide_config["text"],
                        slide_config.get("text_color", "white")
                    )
                elif slide_config["type"] == "chart":
                    self._add_chart_slide(
                        prs,
                        slide_config["chart_path"],
                        image_path
                    )
                elif slide_config["type"] == "insights":
                    self._add_insights_slide(
                        prs,
                        image_path,
                        slide_config["title"],
                        slide_config["points"],
                        slide_config.get("text_color", "white")
                    )
            
            prs.save(pptx_path)
            print(f"PowerPoint created in {time.time() - start_time:.2f} seconds")
            
            return {
                "image_path": image_path,
                "pptx_path": pptx_path,
                "total_time": time.time() - start_time
            }
            
        except Exception as e:
            raise RuntimeError(f"Failed to generate presentation: {str(e)}")

    def _create_translucent_overlay(self, width, height, opacity, is_dark, output_path):
        """Create translucent overlay for slides"""
        if os.path.exists(output_path):
            return

        color = (0, 0, 0, int(255 * opacity)) if is_dark else (255, 255, 255, int(255 * opacity))
        image = Image.new('RGBA', (width, height), color)
        image.save(output_path)

    def _create_slide_with_image(self, prs, image_path, text, text_color='white'):
        """Create a slide with background image and text"""
        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        try:
            # Open the image
            with Image.open(image_path) as img:
                # Calculate aspect ratios
                img_ratio = img.width / img.height
                slide_ratio = slide_width / slide_height

                if img_ratio > slide_ratio:
                    # Image is wider, crop the sides
                    new_width = int(img.height * slide_ratio)
                    left = (img.width - new_width) // 2
                    img = img.crop((left, 0, left + new_width, img.height))
                else:
                    # Image is taller, crop the top and bottom
                    new_height = int(img.width / slide_ratio)
                    top = (img.height - new_height) // 2
                    img = img.crop((0, top, img.width, top + new_height))

                # Convert dimensions to EMU
                img_width = int(slide_width)
                img_height = int(slide_height)

            # Add the image to the slide
            slide.shapes.add_picture(image_path, 0, 0, width=img_width, height=img_height)

            # Create and add the translucent overlay with darker opacity for title slide
            overlay_filename = os.path.join(os.path.dirname(image_path), 
                                          'dark_overlay_title.png' if text_color.lower() == 'white' else 'light_overlay_title.png')
            self._create_translucent_overlay(
                width=img_width // 12700,
                height=img_height // 12700,
                opacity=0.7,  # Even darker overlay (70% opacity)
                is_dark=text_color.lower() == 'white',
                output_path=overlay_filename
            )
            slide.shapes.add_picture(overlay_filename, 0, 0, width=img_width, height=img_height)

            # Add text box - centered on slide
            left = Inches(1)
            top = Inches(2)  # Move down for better positioning
            width = Inches(12)  # Wider text box
            height = Inches(5)
            
            # Calculate center position
            left = (prs.slide_width - width) / 2  # Center horizontally
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center text

            # Split text into lines
            lines = text.split('\n')

            # First line is the title
            p = tf.add_paragraph()
            p.text = lines[0]
            p.font.size = Pt(54)
            p.font.color.rgb = RGBColor(255, 255, 255) if text_color.lower() == 'white' else RGBColor(0, 0, 0)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(20)  # Add space after title

            # Add remaining lines as subtitles
            for line in lines[1:]:
                if line:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(28)
                    p.font.color.rgb = RGBColor(255, 255, 255) if text_color.lower() == 'white' else RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.CENTER
                    p.space_after = Pt(10)  # Add space between subtitle lines

        except Exception as e:
            print(f"An error occurred while creating the slide: {str(e)}")
            traceback.print_exc()

    def _add_chart_slide(self, prs, chart_path, background_image_path):
        """Add a slide with a chart overlay"""
        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        try:
            # 1. Add the background image
            slide.shapes.add_picture(background_image_path, 0, 0, width=slide_width, height=slide_height)

            # 2. Add the white overlay
            overlay_filename = 'white_overlay.png'
            self._create_translucent_overlay(
                width=slide_width // 12700,
                height=slide_height // 12700,
                opacity=0.8,  # 80% white
                is_dark=False,
                output_path=overlay_filename
            )
            slide.shapes.add_picture(overlay_filename, 0, 0, width=slide_width, height=slide_height)

            # 3. Add the chart image on top
            with Image.open(chart_path) as img:
                img_ratio = img.width / img.height
                slide_ratio = slide_width / slide_height

                # Calculate new dimensions to fit 95% of slide width
                new_width = int(slide_width * 0.95)
                new_height = int(new_width / img_ratio)

                # If height exceeds 95% of slide height, scale based on height instead
                if new_height > slide_height * 0.95:
                    new_height = int(slide_height * 0.95)
                    new_width = int(new_height * img_ratio)

                # Add the chart to the slide
                left = (slide_width - new_width) / 2
                top = (slide_height - new_height) / 2
                slide.shapes.add_picture(chart_path, left, top, width=new_width, height=new_height)

        except Exception as e:
            print(f"An error occurred while creating the slide: {str(e)}")
            traceback.print_exc()

    def _add_insights_slide(self, prs, background_image_path, title, points, text_color='white'):
        """Add a slide with bullet points for insights"""
        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        try:
            # 1. Add the background image
            slide.shapes.add_picture(background_image_path, 0, 0, width=slide_width, height=slide_height)

            # 2. Add the overlay
            overlay_filename = os.path.join(os.path.dirname(background_image_path), 
                                          'dark_overlay_insights.png')
            self._create_translucent_overlay(
                width=slide_width // 12700,
                height=slide_height // 12700,
                opacity=0.6,
                is_dark=text_color.lower() == 'white',
                output_path=overlay_filename
            )
            slide.shapes.add_picture(overlay_filename, 0, 0, width=slide_width, height=slide_height)

            # 3. Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = title
            title_para.font.size = Pt(40)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255) if text_color.lower() == 'white' else RGBColor(0, 0, 0)
            title_para.alignment = PP_ALIGN.CENTER

            # 4. Add bullet points with adjusted width and left margin
            left_margin = Inches(1)
            width = Inches(9)
            content_box = slide.shapes.add_textbox(left_margin, Inches(1.5), width, Inches(5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for point in points:
                p = content_frame.add_paragraph()
                
                # Split into bold part and regular part
                if "**" in point:
                    parts = point.split("**: ")
                    bold_text = parts[0].replace("**", "")
                    regular_text = parts[1] if len(parts) > 1 else ""
                    
                    # Add bold part
                    run = p.add_run()
                    run.text = f"• {bold_text}: "
                    run.font.bold = True
                    
                    # Add regular part
                    run = p.add_run()
                    run.text = regular_text
                    run.font.bold = False
                else:
                    p.text = f"• {point}"
                
                # Common formatting for the paragraph
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(255, 255, 255) if text_color.lower() == 'white' else RGBColor(0, 0, 0)
                p.space_after = Pt(14)
                p.level = 0
                p.alignment = PP_ALIGN.LEFT

        except Exception as e:
            print(f"An error occurred while creating the insights slide: {str(e)}")
            traceback.print_exc()